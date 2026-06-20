"""Final-defense deck for the Chest X-ray Pneumonia Detection System.

Main deck = the 15 slides required by the official template (15-minute talk),
redesigned low-text + visual, with per-slide speaker notes. A clearly labelled
Technical Backup / Appendix follows for deep Q&A (not part of the timed talk).

Run: "ai/venv/Scripts/python.exe" documentation/_generators/deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = "D:/Collage/graduation project"
FIG = ROOT + "/documentation/figures/"
SHOTS = ROOT + "/documentation/_verify/shots/"
LOGO = ROOT + "/Zewail Logo/11.png"
DEMO = ROOT + "/Backend/model_assets/demo_output.png"
OUT = ROOT + "/documentation/Defense_Presentation.pptx"

# ---- palette ----
NAVY = RGBColor(0x1F, 0x3B, 0x63)
NAVY2 = RGBColor(0x15, 0x2A, 0x49)
TEAL = RGBColor(0x2A, 0x9D, 0x8F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2B, 0x2F, 0x36)
GREY = RGBColor(0x6A, 0x73, 0x80)
LIGHT = RGBColor(0xED, 0xF2, 0xF8)
BG = RGBColor(0xF7, 0xF9, 0xFC)
LINE = RGBColor(0xD7, 0xDF, 0xEA)
ALERT = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xE7, 0x6F, 0x51)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
FONT = "Calibri"


# ---------------- helpers ----------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, bg)
    return s


def rect(s, l, t, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(l), int(t), int(w), int(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def rrect(s, l, t, w, h, fill, line=None, line_w=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(l), int(t), int(w), int(h))
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def textbox(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4, leading=None):
    tb = s.shapes.add_textbox(int(l), int(t), int(w), int(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        if leading:
            p.line_spacing = leading
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT
    return tb


def bullets(s, l, t, w, h, items, size=18, color=DARK, marker="–", space=10, leading=1.05):
    tb = s.shapes.add_textbox(int(l), int(t), int(w), int(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        bold = False
        if isinstance(it, tuple):
            it, bold = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space); p.line_spacing = leading
        r = p.add_run(); r.text = (marker + "  " + it) if marker else it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT
    return tb


def header(s, title, tag, kicker=None):
    rect(s, 0, 0, SW, Inches(1.12), NAVY)
    rect(s, 0, Inches(1.12), SW, Inches(0.06), TEAL)
    ty = Inches(0.30) if kicker else Inches(0.12)
    if kicker:
        textbox(s, Inches(0.55), Inches(0.14), Inches(10.5), Inches(0.3),
                [(kicker.upper(), 12, RGBColor(0x9F, 0xC0, 0xE0), True)])
    textbox(s, Inches(0.55), ty, Inches(10.6), Inches(0.8),
            [(title, 27, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE if not kicker else MSO_ANCHOR.TOP)
    s.shapes.add_picture(LOGO, int(SW - Inches(1.62)), int(Inches(0.2)), height=int(Inches(0.74)))
    # footer
    textbox(s, Inches(0.45), SH - Inches(0.42), Inches(9.5), Inches(0.35),
            [("Chest X-ray Pneumonia Detection System   |   Team 19   |   CSAI, Zewail City", 9, GREY, False)])
    textbox(s, SW - Inches(2.0), SH - Inches(0.42), Inches(1.6), Inches(0.35),
            [(tag, 9, GREY, False)], align=PP_ALIGN.RIGHT)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def img_fit(s, path, l, t, max_w, max_h, caption=None, cap_color=GREY):
    iw, ih = Image.open(path).size
    r = min(max_w / iw, max_h / ih)
    w, h = int(iw * r), int(ih * r)
    left = int(l + (max_w - w) / 2); top = int(t)
    s.shapes.add_picture(path, left, top, width=w, height=h)
    if caption:
        textbox(s, l, t + h + Inches(0.04), max_w, Inches(0.35),
                [(caption, 11, cap_color, False)], align=PP_ALIGN.CENTER)
    return left, top, w, h


def metric_card(s, l, t, w, h, value, label, vcolor=TEAL):
    rrect(s, l, t, w, h, WHITE, line=LINE, line_w=1.0)
    textbox(s, l, t + Inches(0.16), w, Inches(0.7), [(value, 34, vcolor, True)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, l, t + h - Inches(0.62), w, Inches(0.55), [(label, 12.5, DARK, False)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def card(s, l, t, w, h, title, body, accent=TEAL):
    rrect(s, l, t, w, h, WHITE, line=LINE, line_w=1.0)
    rect(s, l, t + Inches(0.14), Inches(0.08), h - Inches(0.28), accent)
    tb = s.shapes.add_textbox(int(l + Inches(0.26)), int(t + Inches(0.12)), int(w - Inches(0.42)), int(h - Inches(0.24)))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.space_before = Pt(2); r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(11.5); r2.font.color.rgb = DARK; r2.font.name = FONT


def table(s, l, t, w, headers, rows, col_w=None, fsize=12, hfsize=12, row_h=0.4,
          header_fill=NAVY, header_color=WHITE, zebra=True, bold_first=False):
    nrows, ncols = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nrows, ncols, int(l), int(t), int(w), int(Inches(row_h) * nrows)).table
    gt.first_row = False; gt.horz_banding = False
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = int(cw)

    def style(c, text, size, color, bold, fill, align=PP_ALIGN.LEFT):
        c.text = str(text)
        c.fill.solid(); c.fill.fore_color.rgb = fill
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.09); c.margin_right = Inches(0.07)
        c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
        for p in c.text_frame.paragraphs:
            p.alignment = align
            for r in p.runs:
                r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = FONT
    for j, htext in enumerate(headers):
        style(gt.cell(0, j), htext, hfsize, header_color, True, header_fill)
    for i, row in enumerate(rows, start=1):
        fill = LIGHT if (zebra and i % 2 == 0) else WHITE
        for j, val in enumerate(row):
            bold = bold_first and j == 0
            color = NAVY if bold else DARK
            style(gt.cell(i, j), val, fsize, color, bold, fill)
    return gt


def divider(title, subtitle):
    s = slide(NAVY)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, Inches(0.9), Inches(3.05), Inches(0.9), Inches(0.08), TEAL)
    textbox(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(1.0), [(title, 40, WHITE, True)])
    textbox(s, Inches(0.92), Inches(4.25), Inches(11.0), Inches(0.8), [(subtitle, 18, RGBColor(0xBF, 0xD7, 0xEA), False)])
    s.shapes.add_picture(LOGO, int(SW - Inches(2.0)), int(SH - Inches(1.25)), height=int(Inches(0.8)))
    return s


def pipeline(s, l, t, total_w, steps, box_h_in=1.0, color=NAVY, fsize=13):
    """Horizontal block diagram: rounded boxes joined by arrows. steps = list of (title, sub)."""
    n = len(steps)
    arrow_w = int(Inches(0.38))
    box_w = (int(total_w) - arrow_w * (n - 1)) / n
    bh = Inches(box_h_in)
    for i, step in enumerate(steps):
        title, sub = step if isinstance(step, tuple) else (step, None)
        x = int(l + i * (box_w + arrow_w))
        rrect(s, x, t, box_w, bh, color)
        tb = s.shapes.add_textbox(x, int(t), int(box_w), int(bh))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.size = Pt(fsize); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
        if sub:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(1)
            r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(fsize - 3.5)
            r2.font.color.rgb = RGBColor(0xC7, 0xDD, 0xEA); r2.font.name = FONT
        if i < n - 1:
            ax = int(x + box_w)
            textbox(s, ax, int(t), arrow_w, int(bh), [("→", 20, TEAL, True)],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# MAIN DECK  (15 slides, official template order)
# ============================================================

# ---------- 1. Title ----------
s = slide(WHITE)
rect(s, 0, 0, SW, Inches(0.28), NAVY)
rect(s, 0, SH - Inches(0.28), SW, Inches(0.28), TEAL)
s.shapes.add_picture(LOGO, int((SW - Inches(2.7)) / 2), int(Inches(0.55)), width=int(Inches(2.7)))
textbox(s, Inches(1), Inches(2.32), Inches(11.33), Inches(0.45),
        [("Zewail City of Science and Technology", 17, GREY, False)], align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(2.72), Inches(11.93), Inches(1.0),
        [("Chest X-ray Pneumonia Detection System", 39, NAVY, True)], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(3.74), Inches(11.33), Inches(0.45),
        [("An Explainable, Full-Stack Deep-Learning Screening Platform", 16, TEAL, True)], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.32), Inches(11.33), Inches(0.4),
        [("Graduation Project Defense   |   School of CSAI   |   June 2026", 13.5, DARK, False)], align=PP_ALIGN.CENTER)
table(s, Inches(3.07), Inches(4.8), Inches(7.2),
      ["Student", "ID", "Program"],
      [["Moamen Elsayed Elsharkawy", "202202015", "CSAI"],
       ["Habiba Ayman Amin", "202202088", "CSAI"],
       ["Ahmed Gamal Abdelfattah", "202201638", "CSAI"],
       ["Sara Mostafa Ali", "202201305", "CSAI"]],
      col_w=[Inches(4.0), Inches(1.7), Inches(1.5)], fsize=12, hfsize=12, row_h=0.33)
textbox(s, Inches(1), Inches(6.66), Inches(11.33), Inches(0.5),
        [("Supervisor: Prof. Dr. Khaled Mostafa      |      Team Number: 19", 14, NAVY, True)], align=PP_ALIGN.CENTER)
notes(s, "[SPEAKER: Moamen] Good morning. We are Team 19. Our project is an explainable, full-stack pneumonia "
         "screening platform built on the RSNA chest X-ray dataset. Over the next 15 minutes we "
         "walk from the clinical problem, through the models and the system, to honest results, "
         "and finish with a live demo. Supervisor: Prof. Dr. Khaled Mostafa. "
         "Delivery (per the feedback): keep the whole talk to ~15 minutes (up to ~30 with "
         "discussion); every member presents a part and must be able to field any question, AI or "
         "software; have the recorded demo ready as a fallback; present from a local copy, not only "
         "the Drive link; keep slides high-contrast for the projector.")

# ---------- 2. Problem Statement ----------
s = slide(); header(s, "Problem Statement", "2 / 15", kicker="What we are solving")
bullets(s, Inches(0.7), Inches(1.55), Inches(7.5), Inches(5),
        ["Pneumonia is common and dangerous, and the chest X-ray is the first-line test to find it.",
         "The task: given a chest X-ray, decide pneumonia vs normal, localize the region, and report it.",
         "Done manually, this is slow and depends on scarce radiologist expertise.",
         "No accessible, explainable, end-to-end tool does this while keeping the doctor in control.",
         "That gap - automatic X-ray reading + localization + report - is the problem we solve."],
        size=17)
img_fit(s, FIG + "xai_orig.png", Inches(8.5), Inches(1.7), Inches(4.2), Inches(4.4),
        caption="Input: a frontal chest radiograph")
notes(s, "[SPEAKER: Moamen] Note the distinction the panel asked us for: THIS slide is the problem - what the system "
         "must do (read a chest X-ray, decide pneumonia vs normal, localize it, produce a report) and "
         "the gap that no accessible, explainable tool fills it. The why - scarcity, time, cost, human "
         "error - is the next slide, Motivation. Keeping the two separate is exactly the feedback.")

# ---------- 3. Motivation and Objectives ----------
s = slide(); header(s, "Motivation and Objectives", "3 / 15", kicker="Why and what")
textbox(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(0.4), [("Motivation - why it matters", 16, TEAL, True)])
bullets(s, Inches(0.7), Inches(1.82), Inches(11.9), Inches(1.2),
        ["Radiologists are scarce and X-ray volume is high - reads get delayed and doctors are overloaded.",
         "Manual reading is tiring and error-prone; an AI second reader cuts missed cases and saves time.",
         "Existing commercial tools work but are closed and expensive - little fits smaller clinics."], size=14)
textbox(s, Inches(0.7), Inches(3.1), Inches(11.9), Inches(0.4), [("Objectives", 16, TEAL, True)])
ow = Inches(2.92)
objs = [("Evaluate models", "Train and rigorously, leak-free, benchmark detection + localization."),
        ("Ship a product", "Secure multi-user web app: patients, upload, history."),
        ("Explain every call", "A visual heatmap with every single prediction."),
        ("Stay reproducible", "Documented, versioned, maintainable end to end.")]
for i, (tt, bb) in enumerate(objs):
    card(s, Inches(0.7) + i * (ow + Inches(0.12)), Inches(3.5), ow, Inches(2.0), tt, bb)
notes(s, "[SPEAKER: Moamen] Motivation in one line: turn a research model into a safety net a clinician can actually "
         "use. Four objectives anchor the talk - rigorous leak-free evaluation, a real secure "
         "product, explainability on every prediction, and reproducibility. Each later slide maps "
         "back to one of these.")

# ---------- 4. Literature Review / Existing Solutions ----------
s = slide(); header(s, "Literature Review and Existing Solutions", "4 / 15", kicker="Where we build from")
bullets(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.5),
        ["Classification: CheXNet (121-layer DenseNet) reached radiologist-level pneumonia detection; ImageNet backbones transfer well.",
         "Localization: two-stage Faster R-CNN (accurate) vs one-stage YOLO (fast) - both standard on the RSNA challenge (~26,684 X-rays).",
         "Explainability: gradient methods (Grad-CAM, Integrated Gradients, GradientSHAP) and gradient-free (Score-CAM, Eigen-CAM, occlusion)."],
        size=14, space=7)
table(s, Inches(0.7), Inches(3.55), Inches(11.93),
      ["Solution", "Type", "Localization", "Explainability", "Accessibility"],
      [["Lunit INSIGHT CXR", "Commercial", "Yes", "Heatmaps", "Enterprise, paid"],
       ["Aidoc", "Commercial", "Yes", "Limited", "Enterprise, paid"],
       ["qure.ai qXR", "Commercial", "Yes", "Heatmaps", "Enterprise, paid"],
       ["Typical academic model", "Research", "Sometimes", "Rare", "Notebook only"],
       ["This project", "Academic, full-stack", "Yes (boxes)", "6 methods", "Open, self-hostable"]],
      col_w=[Inches(3.0), Inches(2.3), Inches(2.0), Inches(2.1), Inches(2.5)], fsize=12, row_h=0.42)
notes(s, "[SPEAKER: Moamen] The science is settled enough to build on: CNN backbones classify near radiologist level, "
         "and both detector families work on RSNA. Our gap is not a new algorithm - it is the "
         "combination the last row shows: leak-free evaluation across five models, six explainability "
         "methods, delivered as an open, self-hostable product rather than a notebook. Commercial "
         "tools do this but are closed and expensive.")

# ---------- 5. Proposed Solution ----------
s = slide(); header(s, "Proposed Solution", "5 / 15", kicker="The end-to-end workflow")
pipeline(s, Inches(0.55), Inches(1.5), Inches(12.25),
         [("Chest X-ray", "input"), ("Preprocessing", "CLAHE"),
          ("Classification", "normal / pneumonia"), ("Detection", "localize + box"),
          ("Report", "auto-filled")], box_h_in=0.95)
bullets(s, Inches(0.7), Inches(2.95), Inches(7.0), Inches(3.4),
        ["A doctor-facing web app: sign in, manage patients, upload an X-ray, get a reading.",
         "Output: pneumonia decision + confidence, a localized bounding box, and an explainability heatmap.",
         "An automatic templated report summarizes the finding and each detected region with its location.",
         "Five models behind one registry (Faster R-CNN deployed); every analysis stored per patient.",
         "The physician confirms every result - decision support, not autonomy."], size=15, space=9)
img_fit(s, DEMO, Inches(8.15), Inches(3.0), Inches(4.6), Inches(3.35),
        caption="Sample detector output: localized region")
notes(s, "[SPEAKER: Habiba] The supervisor asked us to lead with the workflow as blocks - that is the strip across the "
         "top: an X-ray goes through preprocessing, classification (normal vs pneumonia), detection "
         "(localizing the region), then an auto-filled templated report. The product wraps that "
         "pipeline in a guided clinical web app, stores every analysis per patient, and keeps the "
         "physician as the final decision-maker. Five models sit behind one registry; Faster R-CNN is "
         "the default for its recall.")

# ---------- 6. System Architecture ----------
s = slide(); header(s, "System Architecture", "6 / 15", kicker="Three tiers + the integrated AI model")
img_fit(s, FIG + "architecture.png", Inches(0.6), Inches(1.42), Inches(8.4), Inches(5.6))
for i, (tt, bb) in enumerate([
        ("1. Presentation", "Angular 21 SPA (SSR) - the doctor's browser UI."),
        ("2. Application", "Async FastAPI - REST API, auth, routing."),
        ("3. Data", "MongoDB (Motor) - users, patients, analyses."),
        ("AI model link", "Inference service: registry + warmup; the model output feeds the app + is saved.")]):
    card(s, Inches(9.2), Inches(1.5) + i * Inches(1.36), Inches(3.45), Inches(1.22), tt, bb)
notes(s, "[SPEAKER: Habiba] This slide is the architecture, not the technologies (those are a separate slide, per the "
         "feedback). Three tiers: Angular SSR presentation, async FastAPI application, MongoDB data - "
         "stateless API, all state in the DB. The fourth box is what the panel asked us to show: how "
         "the AI model integrates - a dedicated inference service loads every model once via a registry "
         "and its output flows back into the app and is stored. User flow: doctor opens the browser -> "
         "frontend -> FastAPI -> inference service / model -> result + report -> saved and displayed.")

# ---------- 7. Methodology ----------
s = slide(); header(s, "Methodology", "7 / 15", kicker="Phased AI pipeline + runtime data flow")
img_fit(s, FIG + "ai_pipeline.png", Inches(0.5), Inches(1.45), Inches(12.3), Inches(2.35),
        caption="Research track: an 8-phase pipeline run once per model (prep -> train -> optimize -> retrain -> explain -> evaluate -> demo)")
img_fit(s, FIG + "dataflow.png", Inches(0.5), Inches(4.35), Inches(12.3), Inches(2.2),
        caption="Runtime: upload -> validate -> warmed model -> NMS -> heatmap -> persist + return")
notes(s, "[SPEAKER: Habiba] Two tracks. The research track is an eight-phase pipeline run identically for every model "
         "- the top strip - which is what makes the comparison fair and the numbers reproducible. "
         "Preprocessing applies the same CLAHE contrast to every image regardless of label, which is "
         "how we avoid the leak we hit early on, and the split is patient-wise. The bottom strip is "
         "the live request path. Both are one diagram each so the protocol is auditable.")

# ---------- 8. Main Features ----------
s = slide(); header(s, "Main Features", "8 / 15", kicker="What the doctor gets")
feats = [("Secure accounts", "JWT + bcrypt, per-user data isolation."),
         ("Patient records", "Create, manage, and persist analysis history."),
         ("Upload + preview", "Drag-and-drop X-ray upload with live preview."),
         ("Decision + box", "Pneumonia call, localization, confidence score."),
         ("Explainability", "Grad-CAM and 5 more methods on every prediction."),
         ("Model registry", "Multiple models + auto-generated Swagger API docs.")]
cw, cgap = Inches(3.78), Inches(0.16)
for i, (tt, bb) in enumerate(feats):
    cx = Inches(0.7) + (i % 2) * (cw + cgap)
    cy = Inches(1.55) + (i // 2) * Inches(1.18)
    card(s, cx, cy, cw, Inches(1.04), tt, bb)
img_fit(s, FIG + "ui_result.png", Inches(8.65), Inches(1.55), Inches(4.0), Inches(4.55),
        caption="Result screen: annotated image, confidence, heatmap")
notes(s, "[SPEAKER: Habiba] Six features, but they exist to support one workflow. Security and per-user isolation are "
         "first because this is health data. The clinically important pair is decision-plus-box and "
         "explainability: the doctor sees not just whether, but where, and why. The model registry is "
         "the engineering feature that makes the rest flexible. The screenshot is the real result screen.")

# ---------- 9. Technical Implementation ----------
s = slide(); header(s, "Technical Implementation", "9 / 15", kicker="Stack and engineering")
table(s, Inches(0.7), Inches(1.5), Inches(6.7),
      ["Layer", "Technology"],
      [["Frontend", "Angular 21 (TypeScript), SSR via Express"],
       ["Backend", "FastAPI (Python), Uvicorn ASGI, async I/O"],
       ["Database", "MongoDB with the Motor async driver"],
       ["AI / ML", "PyTorch, torchvision, Ultralytics"],
       ["Auth & Security", "JWT (python-jose), bcrypt (passlib), CORS"],
       ["API tooling", "Postman - design + test the REST API"],
       ["Model storage", "Git LFS for the .pt checkpoints"]],
      col_w=[Inches(1.9), Inches(4.8)], fsize=12, row_h=0.47)
textbox(s, Inches(0.7), Inches(5.35), Inches(6.7), Inches(0.9),
        [("Functional = the features/workflow above.  Non-functional (security, performance, "
          "scalability) improve the existing system - no new feature.", 10.5, GREY, False)])
textbox(s, Inches(7.7), Inches(1.5), Inches(5.0), Inches(0.4), [("Engineering highlights", 15, TEAL, True)])
for i, (tt, bb) in enumerate([
        ("Single-load registry", "Models load once at startup; warmup avoids per-request cost."),
        ("Non-max suppression", "Overlapping detector boxes merged before display."),
        ("Async, non-blocking", "Lifespan startup connects Mongo and warms the model."),
        ("Hardened auth", "Per-request JWT verify + production secret-key guard.")]):
    card(s, Inches(7.7), Inches(1.95) + i * Inches(1.18), Inches(5.0), Inches(1.04), tt, bb)
notes(s, "[SPEAKER: Ahmed] The stack is modern and deliberately boring where it should be. Two engineering decisions "
         "are worth defending: first, the single-load model registry with startup warmup - inference "
         "never pays a load cost mid-request; second, NMS to clean overlapping boxes. The backend is "
         "fully async with a lifespan startup that connects Mongo and warms the model, and auth is "
         "hardened with a guard that refuses to boot in production under the default secret key. "
         "Code listings are in the backup slides.")

# ---------- 10. Testing and Evaluation ----------
s = slide(); header(s, "Testing and Evaluation", "10 / 15", kicker="Two distinct kinds: model vs software")
textbox(s, Inches(0.7), Inches(1.45), Inches(6.0), Inches(0.4), [("Model evaluation", 15, TEAL, True)])
bullets(s, Inches(0.7), Inches(1.9), Inches(6.1), Inches(4.5),
        ["RSNA dataset; DICOM -> PNG; identical CLAHE on every image.",
         "Patient-wise stratified split (seed 42); 20% test (4,135 normal / 1,202 pneumonia).",
         "Per-model metrics: AUC, recall, specificity, F1, confusion matrix; mAP for detection.",
         "Detection: is the region localized correctly (tight box vs too wide / too narrow).",
         "Recall prioritized - a missed pneumonia is the costly error."], size=13.5, space=7)
textbox(s, Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.4), [("Software testing", 15, TEAL, True)])
bullets(s, Inches(7.1), Inches(1.9), Inches(5.5), Inches(4.5),
        ["Unit + integration tests (pytest) on the backend.",
         "API testing with Postman collections on every endpoint.",
         "UI testing of the upload -> result -> history flow.",
         "Per-user access-control tests (a doctor cannot read another's data).",
         "Preflight asset checker + end-to-end inference smoke test."], size=13.5, space=7)
notes(s, "[SPEAKER: Ahmed] The panel was explicit: do not conflate the two. Left is MODEL evaluation - how good the AI "
         "is: the patient-wise split prevents leakage, then AUC, recall, specificity, F1 and mAP per "
         "model, and for detection whether the box actually lands on the finding. Right is SOFTWARE "
         "testing - whether the application works: pytest unit and integration tests, API tests with "
         "Postman, UI testing of the flow, and access-control tests proving per-user isolation. "
         "Different questions, different tools.")

# ---------- 11. Results ----------
s = slide(); header(s, "Results", "11 / 15", kicker="Held-out test set, leak-free")
metric_card(s, Inches(0.7), Inches(1.5), Inches(2.85), Inches(1.35), "0.886", "EfficientNet-B0 AUC\n(best classifier)", TEAL)
metric_card(s, Inches(3.7), Inches(1.5), Inches(2.85), Inches(1.35), "0.812", "Faster R-CNN recall\n(deployed detector)", NAVY)
img_fit(s, FIG + "classification_metrics.png", Inches(0.55), Inches(3.05), Inches(6.2), Inches(3.5))
img_fit(s, FIG + "detection_metrics.png", Inches(6.85), Inches(1.5), Inches(6.0), Inches(4.0))
textbox(s, Inches(6.85), Inches(5.55), Inches(6.0), Inches(1.4),
        [("Recall is prioritized: a missed pneumonia is the costly error.", 13, NAVY, True),
         ("Numbers are honest and leak-free; detection mAP@0.5 sits in the published RSNA range (0.32-0.39).", 11.5, GREY, False)])
notes(s, "[SPEAKER: Ahmed] Headline numbers. Best classifier is EfficientNet-B0 at AUC 0.886 - and notably it is also "
         "the smallest, four million parameters. The deployed detector is Faster R-CNN at recall "
         "0.812, mAP@0.5 0.381. We chose it for recall, not mAP, because missing pneumonia is the "
         "expensive error. These are deliberately honest: our detection mAP lands inside the "
         "published RSNA range of roughly 0.32 to 0.39. If a panelist expects 0.99, that was the "
         "leak - we will explain it in challenges.")

# ---------- 12. Challenges and Lessons Learned ----------
s = slide(); header(s, "Challenges and Lessons Learned", "12 / 15", kicker="What was hard, what we learned")
ch = [("The data leak", "Early ~0.99 AUC came from a patient-level leak + label-correlated preprocessing. Fixed with patient-wise splits and uniform CLAHE.", ALERT),
      ("Class imbalance", "Pneumonia is the minority (~1:3.4). Handled with weighted loss, stratified split, recall-first metrics, train-only augmentation.", TEAL),
      ("Explainable detection", "Standard saliency fails on detectors. Used Eigen-CAM and occlusion sensitivity instead.", TEAL),
      ("Optimization reality", "A cheap PSO/GWO/SA proxy search did NOT beat the tuned baseline - a real, honestly reported finding.", NAVY)]
for i, (tt, bb, ac) in enumerate(ch):
    cx = Inches(0.7) + (i % 2) * Inches(6.1)
    cy = Inches(1.6) + (i // 2) * Inches(2.55)
    card(s, cx, cy, Inches(5.85), Inches(2.3), tt, bb, accent=ac)
notes(s, "[SPEAKER: Ahmed] Four honest lessons. The leak is the one we are proudest of catching: near-perfect early "
         "scores were a patient-level leak compounded by preprocessing that correlated with the "
         "label. Fixing it dropped the numbers and earned our trust. Class imbalance we handle "
         "explicitly - details in backup. Detector explainability needed non-gradient methods. And "
         "the nature-inspired optimization actually lost to the baseline under our cheap proxy budget; "
         "we report that rather than hide it, and we kept the baseline weights.")

# ---------- 13. Conclusion ----------
s = slide(); header(s, "Conclusion", "13 / 15", kicker="What we delivered")
bullets(s, Inches(0.7), Inches(1.65), Inches(8.0), Inches(4.5),
        [("A complete, working, explainable full-stack pneumonia screening system.", True),
         ("Five models benchmarked honestly and leak-free under one protocol.", True),
         ("Secure, multi-user, reproducible, competitive with the literature.", True),
         ("The full path proven: from research model to usable clinical product.", True)],
        size=18, space=14)
metric_card(s, Inches(9.0), Inches(1.65), Inches(3.6), Inches(1.5), "0.886", "Best classifier AUC", TEAL)
metric_card(s, Inches(9.0), Inches(3.35), Inches(3.6), Inches(1.5), "0.812", "Deployed recall", NAVY)
metric_card(s, Inches(9.0), Inches(5.05), Inches(3.6), Inches(1.3), "5", "Models, 1 registry", GOLD)
notes(s, "[SPEAKER: Sarah] To conclude: we set out to turn research models into a usable, honest, explainable clinical "
         "tool, and we did. Five models, one leak-free protocol, a secure multi-user product, numbers "
         "that match the literature instead of overstating it. The contribution is the end-to-end path, "
         "delivered and reproducible.")

# ---------- 14. Future Work ----------
s = slide(); header(s, "Future Work", "14 / 15", kicker="Where it goes next")
fw = [("Containerize + CI/CD", "Docker Compose and a pipeline for one-command deployment."),
      ("Monitoring + drift", "Structured logging, metrics, and model-drift detection."),
      ("Out-of-distribution", "Reject non-chest or low-quality images before inference."),
      ("Multi-finding", "Extend beyond pneumonia; validate on more diverse datasets."),
      ("Active learning", "Turn doctor confirmations into new training labels."),
      ("Clinical validation", "Pursue the regulatory path toward a certified tool.")]
for i, (tt, bb) in enumerate(fw):
    cx = Inches(0.7) + (i % 2) * Inches(6.1)
    cy = Inches(1.6) + (i // 3 if False else (i // 2)) * Inches(1.72)
    card(s, cx, cy, Inches(5.85), Inches(1.5), tt, bb)
notes(s, "[SPEAKER: Sarah] Future work splits into engineering and clinical. Engineering: containerization, CI/CD, and "
         "production monitoring with drift detection - the operational gaps. Safety: out-of-distribution "
         "rejection so the model declines images it should not read. And the longer arc: active learning "
         "from doctor feedback, multi-finding detection, and the clinical-validation and regulatory path "
         "that a real deployment requires.")

# ---------- Live Demo (added per supervisor request) ----------
s = slide(); header(s, "Live Demo", "Live Demo", kicker="The system in action")
bullets(s, Inches(0.7), Inches(1.6), Inches(6.6), Inches(4.8),
        ["Sign in as a doctor and open a patient.",
         "Upload a chest X-ray (drag-and-drop, live preview).",
         "The model runs: pneumonia decision + confidence score.",
         "A localized bounding box + an explainability heatmap.",
         "An auto-generated report, saved to the patient's history."], size=16, space=12)
img_fit(s, FIG + "ui_result.png", Inches(7.7), Inches(1.55), Inches(5.0), Inches(5.0),
        caption="Result screen (shown live)")
notes(s, "[SPEAKER: Sarah] Run the live demo here: sign in, upload an X-ray, and walk through the decision, the "
         "localized box, the heatmap, and the saved templated report. Keep a recorded screen capture "
         "ready as a fallback if the network or backend is unavailable, and run from a local copy, not "
         "only the Drive link (it failed live in the review meeting). Any member should be able to drive it.")

# ---------- 15. Questions ----------
s = slide(NAVY)
rect(s, 0, 0, SW, SH, NAVY)
s.shapes.add_picture(LOGO, int((SW - Inches(2.3)) / 2), int(Inches(1.15)), width=int(Inches(2.3)))
textbox(s, Inches(1), Inches(3.05), Inches(11.33), Inches(1.1), [("Thank You", 46, WHITE, True)], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.2), Inches(11.33), Inches(0.7), [("Questions and Discussion", 23, RGBColor(0xBF, 0xD7, 0xEA), True)], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(5.3), Inches(11.33), Inches(1.0),
        [("Chest X-ray Pneumonia Detection System   |   Team 19", 15, WHITE, False),
         ("School of CSAI, Zewail City of Science and Technology   |   June 2026", 13, RGBColor(0xBF, 0xD7, 0xEA), False)],
        align=PP_ALIGN.CENTER)
notes(s, "[SPEAKER: Sarah] Thank you - we are happy to take questions. Backup slides follow with dataset and imbalance "
         "detail, full metric tables, the explainability gallery, the optimization study, API and "
         "security internals, and individual contributions.")

# ============================================================
# TECHNICAL BACKUP / APPENDIX  (not part of the 15-minute talk)
# ============================================================
divider("Technical Backup", "Supporting material for Q&A - not part of the 15-minute talk")

# B1 Dataset & Class Imbalance
s = slide(); header(s, "Backup: Dataset and Class Imbalance", "Backup", kicker="RSNA Pneumonia Detection Challenge")
bullets(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(2.4),
        ["~26,684 frontal chest radiographs, expert bounding-box labels.",
         "DICOM -> 8-bit PNG; identical CLAHE on every image (no label leak).",
         "Patient-wise stratified split, seed 42; 20% held-out test.",
         "Test split: 4,135 normal (77.5%) / 1,202 pneumonia - ratio ~1 : 3.4."], size=13.5, space=8)
textbox(s, Inches(6.9), Inches(1.5), Inches(5.7), Inches(0.4), [("How we handle the imbalance", 15, TEAL, True)])
for i, (tt, bb) in enumerate([
        ("Class-weighted loss", "Weighted cross-entropy with inverse-frequency weights; missing pneumonia is penalized harder."),
        ("Stratified, patient-wise split", "Keeps the pneumonia:normal ratio across splits and prevents leakage."),
        ("Honest metrics, not accuracy", "Judge on recall, AUC, F1, confusion matrix; recall first."),
        ("Train-only augmentation", "Flip, small rotation, brightness/contrast jitter; lifts the minority class.")]):
    card(s, Inches(6.9), Inches(1.95) + i * Inches(1.16), Inches(5.7), Inches(1.02), tt, bb)
notes(s, "If asked 'how do you stop the model just predicting Normal' - this slide. Imbalance is "
         "real at one to three-point-four. We attack it four ways: inverse-frequency weighted loss, "
         "stratified patient-wise splitting, recall-first metrics instead of accuracy, and train-only "
         "augmentation. The combination is what lifts recall on the minority class, which is the whole "
         "point of a screening tool.")

# B2 Preprocessing & the data-leak fix
s = slide(); header(s, "Backup: Preprocessing and the Data-Leak Fix", "Backup", kicker="Why early scores were too good")
bullets(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(4.6),
        [("Symptom: early models scored ~0.99 AUC - implausibly high for RSNA pneumonia.", True),
         "Cause 1 (preprocessing leak): contrast enhancement (CLAHE) had been applied unevenly, correlated with the label.",
         "Cause 2 (patient leak): images from the same patient appeared in both train and test.",
         ("Fix: apply identical CLAHE to every image regardless of label, and split by patient, not by image.", True),
         "Result: scores dropped to an honest, literature-consistent range - and we now trust the protocol over the number.",
         "Lesson: in medical ML, a number that looks too good is a bug report, not a result."], size=15, space=11)
notes(s, "The leak story in full. Two leaks stacked: preprocessing that correlated with the label, "
         "and patient-level contamination across the split. Both inflate scores in ways that vanish "
         "in the real world. The fix is uniform preprocessing and patient-wise splitting. The reason "
         "we lead with this is that it is the strongest evidence of scientific maturity in the project.")

# B3 Model Zoo
s = slide(); header(s, "Backup: Model Zoo - Five Architectures", "Backup", kicker="Why each model")
table(s, Inches(0.7), Inches(1.55), Inches(11.93),
      ["Model", "Family", "Params", "Input", "Why it is here"],
      [["EfficientNet-B0", "Classification", "4.0 M", "224", "Best AUC at the smallest size - efficiency."],
       ["DenseNet121", "Classification", "7.0 M", "224", "CheXNet lineage; strong feature reuse."],
       ["ResNet50", "Classification", "23.5 M", "224", "Robust baseline; benefited from optimization."],
       ["YOLOv8n", "Detection", "3.0 M", "640", "One-stage, fast; localization baseline."],
       ["Faster R-CNN", "Detection", "41.3 M", "640", "Two-stage; highest recall - deployed default."]],
      col_w=[Inches(2.5), Inches(2.1), Inches(1.3), Inches(1.1), Inches(4.9)], fsize=12.5, row_h=0.6, bold_first=True)
notes(s, "Five models span the trade-off space. Three classifiers from light to heavy, two detectors "
         "one-stage versus two-stage. The point of benchmarking all five under one protocol is to show "
         "the choice is evidence-based: EfficientNet-B0 wins classification at the smallest size, and "
         "Faster R-CNN wins on recall for detection, which is why it is deployed.")

# B4 Training setup
s = slide(); header(s, "Backup: Training Setup", "Backup", kicker="Reproducible protocol")
bullets(s, Inches(0.7), Inches(1.55), Inches(6.0), Inches(4.8),
        ["Hardware: Kaggle GPU (P100 / T4), mixed-precision training.",
         "Classifiers: 224px, ImageNet normalization, transfer learning.",
         "Detectors: 640px; box-aware augmentation.",
         "Optimizer: SGD (momentum 0.9); LinearLR warmup + CosineAnnealing.",
         "Early stopping on the validation metric (patience 3).",
         "Per-model report: architecture, all hyperparameters, full metrics, history."], size=14, space=9)
textbox(s, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.4), [("Faster R-CNN (deployed) hyperparameters", 14, TEAL, True)])
table(s, Inches(7.0), Inches(2.0), Inches(5.6),
      ["Parameter", "Value"],
      [["Base", "ResNet50-FPN (pretrained)"],
       ["Epochs / best epoch", "8 run / epoch 5"],
       ["LR / batch / wd", "0.005 / 4 / 5e-4"],
       ["Scheduler", "LinearLR warmup + Cosine"],
       ["Augmentation", "box-aware hflip + jitter"],
       ["Val mAP@0.5 (best)", "0.423"]],
      col_w=[Inches(2.5), Inches(3.1)], fsize=12, row_h=0.5)
notes(s, "The protocol is identical per model, which is what makes the comparison fair. Classifiers "
         "at 224 with ImageNet transfer, detectors at 640 with box-aware augmentation, SGD with a "
         "warmup-cosine schedule, early stopping on the validation metric. Every run emits a full JSON "
         "report - architecture, hyperparameters, metrics, and the epoch history - so anything on this "
         "slide is traceable to a file.")

# B5 Nature-inspired optimization
s = slide(); header(s, "Backup: Nature-Inspired Optimization (PSO / GWO / SA)", "Backup", kicker="An honest negative result")
bullets(s, Inches(0.7), Inches(1.55), Inches(6.1), Inches(4.8),
        ["We searched hyperparameters with Particle Swarm (PSO), Grey Wolf (GWO), and Simulated Annealing (SA).",
         "Budget was deliberately cheap: population 3, 2 iterations, a 1-epoch proxy per candidate.",
         "The 1-epoch proxy is noisy, so the 'best' params often did not generalize.",
         ("Faster R-CNN: mAP@0.5 0.381 (baseline) -> 0.175 (optimized) = worse.", True),
         "ResNet50 was the exception - its optimized checkpoint did win, so we deployed that one.",
         ("Decision: deploy the better checkpoint per model; report the result honestly.", True)], size=13.5, space=8)
table(s, Inches(7.1), Inches(2.1), Inches(5.5),
      ["Algorithm", "Proxy mAP@0.5"],
      [["PSO", "0.224"], ["GWO", "0.122"], ["SA", "0.079"]],
      col_w=[Inches(2.7), Inches(2.8)], fsize=13, row_h=0.55)
textbox(s, Inches(7.1), Inches(4.2), Inches(5.5), Inches(1.8),
        [("Takeaway", 15, TEAL, True),
         ("A bigger search budget would likely help, but reporting that the cheap search lost is the "
          "scientifically honest result - and it is why the deployed detector uses the baseline weights.", 13, DARK, False)])
notes(s, "This is the slide a research-minded examiner will love. We implemented three metaheuristics "
         "to tune hyperparameters, but under an intentionally cheap budget - population three, two "
         "iterations, one-epoch proxy. That proxy is too noisy, so for Faster R-CNN the 'optimized' "
         "settings actually halved mAP. ResNet50 was the exception and we deployed its optimized "
         "checkpoint. The honest framing: per model, deploy whichever checkpoint is better, and report "
         "that the cheap search did not beat a well-tuned baseline. Future work is simply a larger budget.")

# B6 Full metrics tables
s = slide(); header(s, "Backup: Full Metrics - Held-out Test Set", "Backup", kicker="Leak-free, deployed checkpoints")
textbox(s, Inches(0.7), Inches(1.42), Inches(6.0), Inches(0.35), [("Classification", 14, TEAL, True)])
table(s, Inches(0.7), Inches(1.8), Inches(11.93),
      ["Model", "AUC", "Recall", "Specificity", "F1", "Accuracy", "Params"],
      [["EfficientNet-B0", "0.886", "0.765", "0.830", "0.651", "0.815", "4.0 M"],
       ["DenseNet121", "0.883", "0.785", "0.807", "0.641", "0.802", "7.0 M"],
       ["ResNet50", "0.884", "0.761", "0.824", "0.644", "0.810", "23.5 M"]],
      col_w=[Inches(2.6), Inches(1.35), Inches(1.35), Inches(1.7), Inches(1.35), Inches(1.7), Inches(1.5)],
      fsize=12.5, row_h=0.46, bold_first=True)
textbox(s, Inches(0.7), Inches(4.1), Inches(6.0), Inches(0.35), [("Detection", 14, TEAL, True)])
table(s, Inches(0.7), Inches(4.5), Inches(11.93),
      ["Model", "mAP@0.5", "mAP@[.5:.95]", "Recall", "Precision", "Params"],
      [["Faster R-CNN (deployed)", "0.381", "0.124", "0.812", "-", "41.3 M"],
       ["YOLOv8n", "0.346", "0.138", "0.382", "0.396", "3.0 M"]],
      col_w=[Inches(3.4), Inches(1.7), Inches(2.0), Inches(1.5), Inches(1.6), Inches(1.5)],
      fsize=12.5, row_h=0.5, bold_first=True)
textbox(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.5),
        [("Detection mAP is consistent with the published RSNA range (~0.32-0.39); Faster R-CNN is deployed for its recall.", 12, GREY, False)])
notes(s, "The full table behind the headline. Classifiers are close - all around AUC 0.88 - so we "
         "broke the tie on efficiency and picked EfficientNet-B0. On detection, note the deliberate "
         "trade: YOLO has marginally higher mAP at the strict threshold, but Faster R-CNN more than "
         "doubles recall, 0.81 versus 0.38, which is the metric that matters for screening. Every "
         "number here is from the committed leak-free report files.")

# B7 Confusion matrices
s = slide(); header(s, "Backup: Confusion Matrices", "Backup", kicker="Where the errors fall")
img_fit(s, FIG + "confusion_matrices.png", Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.6))
textbox(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.7),
        [("Errors concentrate in false positives (over-calling), not false negatives - the safe failure mode for screening.", 12.5, NAVY, True)])
notes(s, "Confusion matrices per classifier. The shape that matters: our residual errors are mostly "
         "false positives - the model over-calls pneumonia - rather than false negatives. For a "
         "screening tool that is the failure mode you want, because a flagged normal gets a second "
         "human look, whereas a missed pneumonia goes home. This is recall-first design showing up in "
         "the error distribution.")

# B8 Explainability - classifiers
s = slide(); header(s, "Backup: Explainability Gallery - Classification", "Backup", kicker="Same case, five methods")
gal = [("xai_orig", "Original"), ("xai_gradcam", "Grad-CAM"), ("xai_integrated_gradients", "Integrated Gradients"),
       ("xai_gradient_shap", "GradientSHAP"), ("xai_score_cam", "Score-CAM")]
xw = Inches(2.38)
for i, (fn, lab) in enumerate(gal):
    img_fit(s, FIG + fn + ".png", Inches(0.55) + i * (xw + Inches(0.04)), Inches(1.7), xw, Inches(3.6))
    textbox(s, Inches(0.55) + i * (xw + Inches(0.04)), Inches(5.35), xw, Inches(0.4),
            [(lab, 12, NAVY, True)], align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.8),
        [("Gradient-based (Grad-CAM, Integrated Gradients, GradientSHAP) and gradient-free (Score-CAM) attributions agree on the lung region - evidence the model looks at pathology, not artifacts.", 12.5, DARK, False)])
notes(s, "Explainability on a pneumonia case, five views of the same prediction. The reason we show "
         "multiple methods is robustness: when gradient-based and gradient-free attributions agree on "
         "the same lung region, we have real evidence the model is using pathology, not a chest-tube "
         "or a text marker. This is what lets a doctor challenge the model rather than trust it blindly.")

# B9 Explainability - detection + normal
s = slide(); header(s, "Backup: Explainability - Detection and a Normal Case", "Backup", kicker="Detector-appropriate methods")
det = [("xai_det_box", "Faster R-CNN box"), ("xai_yolo_box", "YOLO box"),
       ("xai_eigencam", "Eigen-CAM"), ("xai_occlusion", "Occlusion")]
for i, (fn, lab) in enumerate(det):
    img_fit(s, FIG + fn + ".png", Inches(0.55) + i * Inches(2.45), Inches(1.7), Inches(2.3), Inches(2.6))
    textbox(s, Inches(0.55) + i * Inches(2.45), Inches(4.35), Inches(2.3), Inches(0.35),
            [(lab, 11.5, NAVY, True)], align=PP_ALIGN.CENTER)
img_fit(s, FIG + "xai_normal_orig.png", Inches(10.5), Inches(1.7), Inches(2.3), Inches(2.6), caption="Normal case")
textbox(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.4),
        [("Standard saliency does not apply to detectors, so we use Eigen-CAM and occlusion sensitivity.", 13, NAVY, True),
         ("On a normal X-ray the detector correctly draws no box and the heatmap stays diffuse - the model abstains rather than inventing a finding.", 12, DARK, False)])
notes(s, "Detectors do not expose a single class gradient the way a classifier does, so Grad-CAM does "
         "not transfer. We use Eigen-CAM on the detector backbone and occlusion sensitivity instead. "
         "The right-hand normal case is important for the panel: on a healthy film the detector draws "
         "no box and the heatmap stays diffuse, so the model abstains rather than hallucinating "
         "pathology - exactly what you want clinically.")

# B10 Inference internals
s = slide(); header(s, "Backup: Inference Internals", "Backup", kicker="From request to result")
bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(4.8),
        ["Model registry: each model declared once (weights, type, class names, thresholds); selected by key.",
         "Single load + warmup at startup - no model is reloaded on a request.",
         "Detection post-processing: confidence threshold, then non-maximum suppression (IoU) to merge overlaps.",
         "Confirmed operating confidence 0.25 (tunable per model); calibrated for recall.",
         "Explainability is best-effort: a heatmap failure never crashes a prediction.",
         "Runs on CPU or GPU; the same checkpoints serve both."], size=15, space=11)
notes(s, "How a request becomes a result. The registry declares each model once; selection is by key, "
         "so swapping the served model is a config change. Everything loads and warms at startup. For "
         "detectors we threshold on confidence then run NMS to merge overlapping boxes. Explainability "
         "is wrapped so a heatmap error degrades gracefully instead of failing the prediction. Same "
         "weights run on CPU or GPU.")

# B11 Backend & API
s = slide(); header(s, "Backup: Backend and API", "Backup", kicker="FastAPI, secured")
table(s, Inches(0.7), Inches(1.55), Inches(7.1),
      ["Endpoint", "Method", "Purpose"],
      [["/api/auth/signup, /login, /me", "POST/GET", "Auth + current user"],
       ["/api/patients", "CRUD", "Patient records"],
       ["/api/xray/upload", "POST", "Upload, infer, persist"],
       ["/api/xray, /{id}", "GET/DEL", "List + manage analyses"],
       ["/api/xray/status, /metadata", "GET", "Model status + metadata"],
       ["/health, /docs", "GET", "Health + Swagger (OpenAPI)"]],
      col_w=[Inches(3.5), Inches(1.5), Inches(2.1)], fsize=11.5, row_h=0.5)
textbox(s, Inches(8.1), Inches(1.5), Inches(4.6), Inches(0.4), [("Security", 15, TEAL, True)])
bullets(s, Inches(8.1), Inches(1.95), Inches(4.6), Inches(4.5),
        ["JWT bearer auth, verified per request.",
         "bcrypt password hashing (per-password salt).",
         "Per-user authorization - scoped DB queries.",
         "Strict input validation (Pydantic).",
         "CORS allow-list; production secret-key guard.",
         "Auto-generated OpenAPI docs at /docs."], size=13, space=9)
notes(s, "The API surface is small and RESTful, and the OpenAPI docs at slash-docs are auto-generated "
         "so the contract is always live. On security: JWT verified on every protected request, bcrypt "
         "for passwords, and - the part examiners probe - authorization is enforced at the query level, "
         "so a doctor's queries are scoped to their own data. The production guard refuses to boot "
         "under the default secret key.")

# B12 Data model
s = slide(); header(s, "Backup: Data Model", "Backup", kicker="Three MongoDB collections")
img_fit(s, FIG + "erd.png", Inches(0.8), Inches(1.5), Inches(8.0), Inches(5.2))
bullets(s, Inches(9.0), Inches(1.8), Inches(3.7), Inches(4.5),
        ["users -> patients -> xray_analyses (1-to-many).",
         "Unique indexes on email and identifiers.",
         "Compound index (owner, created_at) for fast history.",
         "Every analysis links to its patient and owning doctor.",
         "Document model fits nested, evolving analysis records."], size=13, space=10)
notes(s, "The data model is three collections - users own patients own analyses. Two indexing "
         "decisions matter: unique indexes prevent duplicate accounts and records, and a compound "
         "index on owner plus creation time makes the history view fast and keeps queries naturally "
         "scoped per doctor. We chose a document store because an analysis record is nested and its "
         "shape evolves.")

# B13 Frontend
s = slide(); header(s, "Backup: Frontend (Angular 21)", "Backup", kicker="SSR, guarded, responsive")
bullets(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(4.8),
        ["Angular 21 single-page app with server-side rendering (Express).",
         "Route guards: protected pages require auth; signed-in users skip public pages.",
         "Central state + a single API service layer (typed HTTP).",
         "Guided linear workflow: dashboard -> upload -> processing -> result.",
         "Responsive - works on mobile (shots at right)."], size=14, space=10)
img_fit(s, SHOTS + "m2_upload.png", Inches(7.1), Inches(1.55), Inches(2.5), Inches(5.1))
img_fit(s, SHOTS + "m3_records.png", Inches(9.9), Inches(1.55), Inches(2.5), Inches(5.1))
notes(s, "The frontend is Angular 21 with SSR for fast first paint. Two engineering points: route "
         "guards enforce auth at navigation time, and all HTTP goes through one typed API service so "
         "the data layer is centralized. The UX is deliberately one linear path - dashboard, upload, "
         "processing, result - so a busy clinician needs no training. It is responsive; these are the "
         "real mobile screens.")

# B14 Deployment & reproducibility
s = slide(); header(s, "Backup: Deployment and Reproducibility", "Backup", kicker="How to run it")
bullets(s, Inches(0.7), Inches(1.6), Inches(8.2), Inches(4.8),
        ["Backend: Python venv, requirements.txt, .env (secret key + Mongo URL), run via uvicorn.",
         "Frontend: npm install, npm start (dev) or npm run build (prod).",
         "Model weights pulled via Git LFS on clone - drop-in .pt files, no code change.",
         "Stateless API: scales horizontally behind a load balancer; DB holds all state.",
         "Every training run is reproducible from its committed report JSON (seed 42).",
         "Full source, assets, thesis, and deck in the GitHub repository."], size=14, space=10)
img_fit(s, FIG + "github_qr.png", Inches(9.6), Inches(1.7), Inches(3.0), Inches(3.0), caption="Project repository")
notes(s, "Reproducibility was a design goal, not an afterthought. The backend is a standard venv plus "
         "env file; the frontend is npm; model weights come down through Git LFS so they are versioned "
         "with the code and drop in without changes. The API is stateless, so it scales horizontally. "
         "And because seed and hyperparameters are committed per run, any number in this deck can be "
         "regenerated. QR links the repo.")

# B15 Individual contributions
s = slide(); header(s, "Backup: Individual Contributions", "Backup", kicker="Team 19")
table(s, Inches(0.7), Inches(1.6), Inches(11.93),
      ["Member", "Track", "Main responsibilities"],
      [["Moamen Elsayed Elsharkawy", "AI Research", "Data pipeline + leak fix, model training, evaluation, explainability, integration lead."],
       ["Habiba Ayman Amin", "AI Research", "Model training/experiments, metrics + confusion matrices, optimization study, XAI."],
       ["Ahmed Gamal Abdelfattah", "Full-stack", "FastAPI backend, MongoDB data layer, inference service wiring, API + auth."],
       ["Sara Mostafa Ali", "Full-stack", "Angular frontend, UI/UX, upload + result + history views, integration testing."]],
      col_w=[Inches(3.4), Inches(1.9), Inches(6.6)], fsize=12.5, row_h=0.78, bold_first=True)
textbox(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.6),
        [("AI track (Moamen, Habiba): data, models, evaluation, explainability.  Full-stack track (Ahmed, Sara): frontend, backend, database, integration.  Individual contributions are reflected in the Git commit history.", 12, GREY, False)])
notes(s, "Required by the template. The team splits into an AI track and a full-stack track. Moamen "
         "and Habiba owned the data pipeline, training, evaluation, and explainability; Ahmed and Sara "
         "owned the Angular frontend, the FastAPI backend, the database, and integration. The Git "
         "history reflects the per-person split. (Confirm the exact wording with the team before the defense.)")

# B16 References
s = slide(); header(s, "Backup: Key References", "Backup", kicker="Selected")
bullets(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(5.2),
        ["Rajpurkar et al. CheXNet: Radiologist-Level Pneumonia Detection on Chest X-Rays with Deep Learning. 2017.",
         "Radiological Society of North America. RSNA Pneumonia Detection Challenge. Kaggle, 2018.",
         "Ren, He, Girshick, Sun. Faster R-CNN: Towards Real-Time Object Detection with Region Proposal Networks. NeurIPS 2015.",
         "Jocher et al. Ultralytics YOLOv8. 2023.",
         "Selvaraju et al. Grad-CAM: Visual Explanations from Deep Networks via Gradient-based Localization. ICCV 2017.",
         "Sundararajan, Taly, Yan. Axiomatic Attribution for Deep Networks (Integrated Gradients). ICML 2017.",
         "Tan, Le. EfficientNet: Rethinking Model Scaling for CNNs. ICML 2019.",
         "Huang et al. Densely Connected Convolutional Networks (DenseNet). CVPR 2017."], size=13, space=8, marker=None)
notes(s, "Selected references - the full IEEE list is in the thesis. These cover the backbones "
         "(DenseNet, EfficientNet, ResNet via CheXNet), both detectors, the dataset, and the "
         "explainability methods we used.")

prs.save(OUT)
print("WROTE", OUT, "-", len(prs.slides._sldIdLst), "slides")
